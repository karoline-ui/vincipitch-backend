"""
Serviço de exportação de documentos (PDF, XLSX, DOCX, PPTX).
Gera relatórios profissionais com gráficos embutidos.
"""

import io
import logging
from typing import Dict, List, Optional, Any
from datetime import datetime
from pathlib import Path

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm, mm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    Image, PageBreak, ListFlowable, ListItem, HRFlowable, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY, TA_RIGHT
from reportlab.graphics.shapes import Drawing, Line, Rect
from reportlab.graphics import renderPDF

from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.drawing.image import Image as XLImage
from openpyxl.chart import RadarChart, Reference, BarChart

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE

from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTXPt
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# Cores da marca
CORES = {
    "verde": colors.HexColor('#00E676'),
    "verde_escuro": colors.HexColor('#00C853'),
    "roxo": colors.HexColor('#4A148C'),
    "roxo_claro": colors.HexColor('#7B1FA2'),
    "azul": colors.HexColor('#1976D2'),
    "cinza_escuro": colors.HexColor('#1E293B'),
    "cinza_medio": colors.HexColor('#475569'),
    "cinza_claro": colors.HexColor('#94A3B8'),
    "fundo_claro": colors.HexColor('#F8FAFC'),
    "borda": colors.HexColor('#E2E8F0'),
    "vermelho": colors.HexColor('#EF4444'),
    "amarelo": colors.HexColor('#F59E0B'),
    "branco": colors.white,
}

CRITERIOS_LABELS = {
    "sumario_executivo": "Sumário Executivo",
    "proposta_valor": "Proposta de Valor",
    "concorrencia": "Concorrência",
    "mercado_alvo": "Mercado Alvo",
    "canais_distribuicao": "Canais de Distribuição",
    "relacionamento_clientes": "Relacionamento com Clientes",
    "fontes_receita": "Fontes de Receita",
    "recursos_principais": "Recursos Principais",
    "atividades_chave": "Atividades-Chave",
    "parceiros": "Parceiros Estratégicos",
    "estrutura_custos": "Estrutura de Custos",
    "referencias_indicacao": "Referências e Indicações"
}

# Mapeamento das justificativas
JUSTIFICATIVA_MAP = {
    'sumario_executivo': 'justificativa_sumario',
    'proposta_valor': 'justificativa_proposta',
    'concorrencia': 'justificativa_concorrencia',
    'mercado_alvo': 'justificativa_mercado',
    'canais_distribuicao': 'justificativa_canais',
    'relacionamento_clientes': 'justificativa_relacionamento',
    'fontes_receita': 'justificativa_receita',
    'recursos_principais': 'justificativa_recursos',
    'atividades_chave': 'justificativa_atividades',
    'parceiros': 'justificativa_parceiros',
    'estrutura_custos': 'justificativa_custos',
    'referencias_indicacao': 'justificativa_referencias',
}


class DocumentExporter:
    """Exportador de documentos para relatórios de análise."""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self._configurar_estilos()
    
    def _configurar_estilos(self):
        """Configura estilos personalizados para PDF."""
        
        # Título principal da capa
        self.styles.add(ParagraphStyle(
            name='TituloCapa',
            parent=self.styles['Heading1'],
            fontSize=32,
            textColor=CORES['roxo'],
            spaceAfter=10,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Subtítulo da capa
        self.styles.add(ParagraphStyle(
            name='SubtituloCapa',
            parent=self.styles['Normal'],
            fontSize=14,
            textColor=CORES['cinza_medio'],
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica'
        ))
        
        # Nome da empresa
        self.styles.add(ParagraphStyle(
            name='NomeEmpresa',
            parent=self.styles['Heading1'],
            fontSize=28,
            textColor=CORES['verde_escuro'],
            spaceBefore=20,
            spaceAfter=10,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Título de seção
        self.styles.add(ParagraphStyle(
            name='TituloSecao',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=CORES['roxo'],
            spaceBefore=20,
            spaceAfter=12,
            fontName='Helvetica-Bold',
            borderPadding=5
        ))
        
        # Subtítulo de seção
        self.styles.add(ParagraphStyle(
            name='SubtituloSecao',
            parent=self.styles['Heading3'],
            fontSize=13,
            textColor=CORES['verde_escuro'],
            spaceBefore=15,
            spaceAfter=8,
            fontName='Helvetica-Bold'
        ))
        
        # Corpo de texto
        self.styles.add(ParagraphStyle(
            name='Corpo',
            parent=self.styles['Normal'],
            fontSize=10,
            textColor=CORES['cinza_escuro'],
            leading=14,
            alignment=TA_JUSTIFY,
            spaceAfter=8
        ))
        
        # Texto pequeno
        self.styles.add(ParagraphStyle(
            name='TextoPequeno',
            parent=self.styles['Normal'],
            fontSize=9,
            textColor=CORES['cinza_medio'],
            leading=12
        ))
        
        # Nota alta (verde)
        self.styles.add(ParagraphStyle(
            name='NotaAlta',
            parent=self.styles['Normal'],
            fontSize=36,
            textColor=CORES['verde_escuro'],
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Nota baixa (vermelho)
        self.styles.add(ParagraphStyle(
            name='NotaBaixa',
            parent=self.styles['Normal'],
            fontSize=36,
            textColor=CORES['vermelho'],
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Nota média (amarelo)
        self.styles.add(ParagraphStyle(
            name='NotaMedia',
            parent=self.styles['Normal'],
            fontSize=36,
            textColor=CORES['amarelo'],
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Lista com bullet
        self.styles.add(ParagraphStyle(
            name='ListaItem',
            parent=self.styles['Normal'],
            fontSize=10,
            textColor=CORES['cinza_escuro'],
            leading=14,
            leftIndent=15,
            spaceAfter=4
        ))
    
    def _criar_linha_decorativa(self, largura=16*cm, cor=None):
        """Cria uma linha decorativa horizontal."""
        if cor is None:
            cor = CORES['roxo']
        return HRFlowable(
            width=largura,
            thickness=2,
            color=cor,
            spaceBefore=10,
            spaceAfter=10,
            hAlign='CENTER'
        )
    
    def _criar_linha_fina(self, largura=16*cm):
        """Cria uma linha fina cinza."""
        return HRFlowable(
            width=largura,
            thickness=0.5,
            color=CORES['borda'],
            spaceBefore=5,
            spaceAfter=5,
            hAlign='CENTER'
        )
    
    def _get_cor_nota(self, nota: float) -> colors.Color:
        """Retorna a cor baseada na nota."""
        if nota >= 3.0:
            return CORES['verde_escuro']
        elif nota >= 2.0:
            return CORES['amarelo']
        else:
            return CORES['vermelho']
    
    def _get_estilo_nota(self, nota: float) -> str:
        """Retorna o estilo baseado na nota."""
        if nota >= 3.0:
            return 'NotaAlta'
        elif nota >= 2.0:
            return 'NotaMedia'
        else:
            return 'NotaBaixa'
    
    # =========================================================================
    # EXPORTAR PDF PROFISSIONAL
    # =========================================================================
    
    async def exportar_pdf(
        self,
        analise: Dict[str, Any],
        empresa: Dict[str, Any],
        grafico_radar: Optional[bytes] = None,
        diagnostico: Optional[Dict[str, Any]] = None
    ) -> bytes:
        """
        Exporta análise completa para PDF com design profissional.
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            rightMargin=2*cm,
            leftMargin=2*cm,
            topMargin=1.5*cm,
            bottomMargin=1.5*cm
        )
        
        story = []
        nota_final = analise.get('nota_final', 0)
        
        # ═══════════════════════════════════════════════════════════════════
        # CAPA
        # ═══════════════════════════════════════════════════════════════════
        
        story.append(Spacer(1, 2*cm))
        
        # Título
        story.append(Paragraph("VinciPitch.AI", self.styles['TituloCapa']))
        story.append(Paragraph("Relatório de Análise de Startup", self.styles['SubtituloCapa']))
        
        story.append(Spacer(1, 0.5*cm))
        story.append(self._criar_linha_decorativa(12*cm, CORES['roxo']))
        story.append(Spacer(1, 1*cm))
        
        # Nome da empresa
        story.append(Paragraph(empresa.get('nome', 'N/A'), self.styles['NomeEmpresa']))
        
        # Setor
        setor = empresa.get('setor', 'N/A')
        setor_formatado = setor.replace('_', ' ').title() if setor else 'N/A'
        story.append(Paragraph(f"Setor: {setor_formatado}", self.styles['SubtituloCapa']))
        
        story.append(Spacer(1, 1.5*cm))
        
        # Card de nota final (tabela para simular card)
        nota_cor = self._get_cor_nota(nota_final)
        nota_data = [
            [Paragraph(f'<font size="28"><b>{nota_final:.2f}</b></font>', 
                      ParagraphStyle('NotaCapa', alignment=TA_CENTER, textColor=CORES['cinza_escuro']))],
            [Paragraph('de 4.00', ParagraphStyle('NotaSubCapa', alignment=TA_CENTER, 
                      textColor=CORES['cinza_medio'], fontSize=10))]
        ]
        nota_table = Table(nota_data, colWidths=[6*cm])
        nota_table.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('BACKGROUND', (0, 0), (-1, -1), CORES['fundo_claro']),
            ('BOX', (0, 0), (-1, -1), 2, nota_cor),
            ('TOPPADDING', (0, 0), (-1, -1), 15),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 15),
        ]))
        nota_table.hAlign = 'CENTER'
        story.append(nota_table)
        
        story.append(Paragraph("Nota Final", self.styles['SubtituloCapa']))
        
        story.append(Spacer(1, 2*cm))
        
        # Data
        story.append(Paragraph(
            f"Gerado em: {datetime.now().strftime('%d/%m/%Y às %H:%M')}",
            ParagraphStyle('DataCapa', alignment=TA_CENTER, textColor=CORES['cinza_claro'], fontSize=10)
        ))
        
        story.append(PageBreak())
        
        # ═══════════════════════════════════════════════════════════════════
        # SUMÁRIO EXECUTIVO
        # ═══════════════════════════════════════════════════════════════════
        
        story.append(Paragraph("Sumário Executivo", self.styles['TituloSecao']))
        story.append(self._criar_linha_fina(16*cm))
        
        resumo = analise.get('resumo_executivo') or (diagnostico.get('resumo_executivo') if diagnostico else None) or 'Análise não disponível.'
        story.append(Paragraph(resumo, self.styles['Corpo']))
        
        story.append(Spacer(1, 0.8*cm))
        
        # ═══════════════════════════════════════════════════════════════════
        # CLASSIFICAÇÕES
        # ═══════════════════════════════════════════════════════════════════
        
        story.append(Paragraph("Classificações", self.styles['TituloSecao']))
        story.append(self._criar_linha_fina(16*cm))
        
        potencial = analise.get('classificacao_potencial', 'N/A')
        risco = analise.get('classificacao_risco', 'N/A')
        recomendacao = analise.get('recomendacao_investimento', 'N/A')
        
        potencial_fmt = potencial.replace('_', ' ').title() if potencial else 'N/A'
        risco_fmt = risco.replace('_', ' ').title() if risco else 'N/A'
        recomendacao_fmt = recomendacao.replace('_', ' ').title() if recomendacao else 'N/A'
        
        classificacoes_data = [
            ['Critério', 'Classificação'],
            ['Potencial de Investimento', potencial_fmt],
            ['Nível de Risco', risco_fmt],
            ['Recomendação', recomendacao_fmt]
        ]
        
        table = Table(classificacoes_data, colWidths=[8*cm, 8*cm])
        table.setStyle(TableStyle([
            # Cabeçalho - roxo com texto BRANCO
            ('BACKGROUND', (0, 0), (-1, 0), CORES['roxo']),
            ('TEXTCOLOR', (0, 0), (-1, 0), CORES['branco']),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
            ('TOPPADDING', (0, 0), (-1, 0), 10),
            # Corpo
            ('BACKGROUND', (0, 1), (-1, -1), CORES['fundo_claro']),
            ('TEXTCOLOR', (0, 1), (-1, -1), CORES['cinza_escuro']),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('TOPPADDING', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
            # Geral
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('GRID', (0, 0), (-1, -1), 0.5, CORES['borda']),
        ]))
        table.hAlign = 'CENTER'
        story.append(table)
        
        story.append(Spacer(1, 0.8*cm))
        
        # ═══════════════════════════════════════════════════════════════════
        # GRÁFICO RADAR
        # ═══════════════════════════════════════════════════════════════════
        
        if grafico_radar:
            story.append(Paragraph("Análise por Critérios", self.styles['TituloSecao']))
            story.append(self._criar_linha_fina(16*cm))
            
            img = Image(io.BytesIO(grafico_radar), width=13*cm, height=13*cm)
            img.hAlign = 'CENTER'
            story.append(img)
        
        story.append(PageBreak())
        
        # ═══════════════════════════════════════════════════════════════════
        # NOTAS DETALHADAS
        # ═══════════════════════════════════════════════════════════════════
        
        story.append(Paragraph("Notas Detalhadas por Critério", self.styles['TituloSecao']))
        story.append(self._criar_linha_fina(16*cm))
        
        # Estilo para células
        cell_style = ParagraphStyle(
            'CellStyle',
            parent=self.styles['Normal'],
            fontSize=9,
            textColor=CORES['cinza_escuro'],
            leading=11
        )
        
        # Cabeçalho da tabela
        header = ['Critério', 'Nota', 'Justificativa']
        notas_data = [header]
        
        for criterio_key, criterio_label in CRITERIOS_LABELS.items():
            nota = analise.get(f'nota_{criterio_key}', 0)
            just_key = JUSTIFICATIVA_MAP.get(criterio_key, f'justificativa_{criterio_key}')
            justificativa_texto = analise.get(just_key, '') or 'N/A'
            
            # Nota formatada com cor
            cor_nota = self._get_cor_nota(nota)
            nota_str = f'{nota:.1f}'
            
            just_para = Paragraph(justificativa_texto, cell_style)
            
            notas_data.append([criterio_label, nota_str, just_para])
        
        # Tabela
        table = Table(notas_data, colWidths=[4*cm, 1.2*cm, 11.3*cm])
        
        # Estilos da tabela
        table_style = [
            # Cabeçalho - roxo com texto BRANCO
            ('BACKGROUND', (0, 0), (-1, 0), CORES['roxo']),
            ('TEXTCOLOR', (0, 0), (-1, 0), CORES['branco']),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('TOPPADDING', (0, 0), (-1, 0), 8),
            # Corpo
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('TOPPADDING', (0, 1), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
            ('LEFTPADDING', (0, 0), (-1, -1), 5),
            ('RIGHTPADDING', (0, 0), (-1, -1), 5),
            # Alinhamento
            ('ALIGN', (0, 0), (0, -1), 'LEFT'),
            ('ALIGN', (1, 0), (1, -1), 'CENTER'),
            ('ALIGN', (2, 0), (2, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            # Bordas
            ('GRID', (0, 0), (-1, -1), 0.5, CORES['borda']),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [CORES['branco'], CORES['fundo_claro']]),
        ]
        
        # Colorir notas individualmente
        for i, (criterio_key, _) in enumerate(CRITERIOS_LABELS.items(), 1):
            nota = analise.get(f'nota_{criterio_key}', 0)
            cor = self._get_cor_nota(nota)
            table_style.append(('TEXTCOLOR', (1, i), (1, i), cor))
            table_style.append(('FONTNAME', (1, i), (1, i), 'Helvetica-Bold'))
        
        table.setStyle(TableStyle(table_style))
        table.hAlign = 'CENTER'
        story.append(table)
        
        story.append(PageBreak())
        
        # ═══════════════════════════════════════════════════════════════════
        # DIAGNÓSTICO SWOT
        # ═══════════════════════════════════════════════════════════════════
        
        if diagnostico:
            story.append(Paragraph("Diagnóstico Estratégico", self.styles['TituloSecao']))
            story.append(self._criar_linha_fina(16*cm))
            
            # Pontos Fortes
            if diagnostico.get('pontos_fortes'):
                story.append(Paragraph("Pontos Fortes", self.styles['SubtituloSecao']))
                for ponto in diagnostico.get('pontos_fortes', [])[:5]:
                    story.append(Paragraph(f"• {ponto}", self.styles['ListaItem']))
                story.append(Spacer(1, 0.4*cm))
            
            # Pontos Fracos
            if diagnostico.get('pontos_fracos'):
                story.append(Paragraph("Pontos de Atenção", self.styles['SubtituloSecao']))
                for ponto in diagnostico.get('pontos_fracos', [])[:5]:
                    story.append(Paragraph(f"• {ponto}", self.styles['ListaItem']))
                story.append(Spacer(1, 0.4*cm))
            
            # Oportunidades
            if diagnostico.get('oportunidades'):
                story.append(Paragraph("Oportunidades", self.styles['SubtituloSecao']))
                for ponto in diagnostico.get('oportunidades', [])[:5]:
                    story.append(Paragraph(f"• {ponto}", self.styles['ListaItem']))
                story.append(Spacer(1, 0.4*cm))
            
            # Ameaças
            if diagnostico.get('ameacas'):
                story.append(Paragraph("Ameaças", self.styles['SubtituloSecao']))
                for ponto in diagnostico.get('ameacas', [])[:5]:
                    story.append(Paragraph(f"• {ponto}", self.styles['ListaItem']))
                story.append(Spacer(1, 0.4*cm))
            
            story.append(self._criar_linha_fina())
            
            # Recomendações
            if diagnostico.get('recomendacoes'):
                story.append(Paragraph("Recomendações", self.styles['SubtituloSecao']))
                for i, rec in enumerate(diagnostico.get('recomendacoes', [])[:5], 1):
                    story.append(Paragraph(f"{i}. {rec}", self.styles['ListaItem']))
                story.append(Spacer(1, 0.4*cm))
            
            # Próximos Passos
            if diagnostico.get('proximos_passos'):
                story.append(Paragraph("Próximos Passos", self.styles['SubtituloSecao']))
                for i, passo in enumerate(diagnostico.get('proximos_passos', [])[:5], 1):
                    story.append(Paragraph(f"{i}. {passo}", self.styles['ListaItem']))
        
        # ═══════════════════════════════════════════════════════════════════
        # RODAPÉ FINAL (centralizado)
        # ═══════════════════════════════════════════════════════════════════
        
        story.append(Spacer(1, 1*cm))
        story.append(self._criar_linha_fina(10*cm))
        story.append(Paragraph(
            "Relatório gerado automaticamente por VinciPitch.AI",
            ParagraphStyle('Rodape', alignment=TA_CENTER, textColor=CORES['cinza_claro'], fontSize=9)
        ))
        story.append(Paragraph(
            "© 2026 - Plataforma de Análise de Startups",
            ParagraphStyle('Rodape2', alignment=TA_CENTER, textColor=CORES['cinza_claro'], fontSize=8)
        ))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    
    # =========================================================================
    # EXPORTAR XLSX
    # =========================================================================
    
    async def exportar_xlsx(
        self,
        analises: List[Dict[str, Any]],
        empresas: List[Dict[str, Any]],
        titulo: str = "Ranking VinciPitch"
    ) -> bytes:
        """
        Exporta ranking para planilha Excel.
        
        Args:
            analises: Lista de análises
            empresas: Lista de empresas correspondentes
            titulo: Título da planilha
            
        Returns:
            XLSX em bytes
        """
        wb = Workbook()
        ws = wb.active
        ws.title = "Ranking"
        
        # Estilos
        header_fill = PatternFill(start_color="4A148C", end_color="4A148C", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=11)
        verde_fill = PatternFill(start_color="00E676", end_color="00E676", fill_type="solid")
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        # Cabeçalho
        headers = [
            "Posição", "Empresa", "Setor", "Nota Final", "Potencial", "Risco",
            "Sumário", "Proposta Valor", "Concorrência", "Mercado",
            "Canais", "Relacionamento", "Receita", "Recursos",
            "Atividades", "Parceiros", "Custos", "Referências"
        ]
        
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center', vertical='center')
            cell.border = thin_border
        
        # Dados
        empresas_dict = {e['id']: e for e in empresas}
        
        for row, analise in enumerate(analises, 2):
            empresa = empresas_dict.get(analise.get('empresa_id'), {})
            
            data = [
                row - 1,  # Posição
                empresa.get('nome', 'N/A'),
                empresa.get('setor', 'N/A').title(),
                analise.get('nota_final', 0),
                analise.get('classificacao_potencial', 'N/A'),
                analise.get('classificacao_risco', 'N/A'),
                analise.get('nota_sumario_executivo', 0),
                analise.get('nota_proposta_valor', 0),
                analise.get('nota_concorrencia', 0),
                analise.get('nota_mercado_alvo', 0),
                analise.get('nota_canais_distribuicao', 0),
                analise.get('nota_relacionamento_clientes', 0),
                analise.get('nota_fontes_receita', 0),
                analise.get('nota_recursos_principais', 0),
                analise.get('nota_atividades_chave', 0),
                analise.get('nota_parceiros', 0),
                analise.get('nota_estrutura_custos', 0),
                analise.get('nota_referencias_indicacao', 0)
            ]
            
            for col, value in enumerate(data, 1):
                cell = ws.cell(row=row, column=col, value=value)
                cell.border = thin_border
                cell.alignment = Alignment(horizontal='center' if col != 2 else 'left')
                
                # Formatação condicional para notas
                if col >= 4 and isinstance(value, (int, float)):
                    if value >= 3.0:
                        cell.fill = PatternFill(start_color="00E676", fill_type="solid")
                    elif value >= 2.0:
                        cell.fill = PatternFill(start_color="FACC15", fill_type="solid")
                    elif value < 1.5:
                        cell.fill = PatternFill(start_color="EF4444", fill_type="solid")
        
        # Ajusta largura das colunas
        column_widths = [8, 30, 15, 12, 12, 12] + [10] * 12
        for i, width in enumerate(column_widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = width
        
        # Congela primeira linha
        ws.freeze_panes = 'A2'
        
        # Adiciona filtros
        ws.auto_filter.ref = ws.dimensions
        
        # ===== ABA DE ESTATÍSTICAS =====
        ws_stats = wb.create_sheet("Estatísticas")
        
        # Calcula estatísticas
        notas = [a.get('nota_final', 0) for a in analises]
        
        stats = [
            ["Métrica", "Valor"],
            ["Total de Empresas", len(analises)],
            ["Média Geral", f"{sum(notas)/len(notas):.2f}" if notas else "0"],
            ["Nota Máxima", f"{max(notas):.2f}" if notas else "0"],
            ["Nota Mínima", f"{min(notas):.2f}" if notas else "0"],
            ["Empresas ≥ 3.0", sum(1 for n in notas if n >= 3.0)],
            ["Empresas ≥ 2.0", sum(1 for n in notas if n >= 2.0)],
            ["Empresas < 2.0", sum(1 for n in notas if n < 2.0)]
        ]
        
        for row, (metrica, valor) in enumerate(stats, 1):
            ws_stats.cell(row=row, column=1, value=metrica)
            ws_stats.cell(row=row, column=2, value=valor)
            if row == 1:
                ws_stats.cell(row=row, column=1).fill = header_fill
                ws_stats.cell(row=row, column=1).font = header_font
                ws_stats.cell(row=row, column=2).fill = header_fill
                ws_stats.cell(row=row, column=2).font = header_font
        
        ws_stats.column_dimensions['A'].width = 20
        ws_stats.column_dimensions['B'].width = 15
        
        # Salva em bytes
        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        
        return buffer.read()
    
    # =========================================================================
    # EXPORTAR DOCX
    # =========================================================================
    
    async def exportar_docx(
        self,
        analise: Dict[str, Any],
        empresa: Dict[str, Any],
        grafico_radar: Optional[bytes] = None,
        diagnostico: Optional[Dict[str, Any]] = None
    ) -> bytes:
        """
        Exporta análise para documento Word.
        
        Args:
            analise: Dados da análise
            empresa: Dados da empresa
            grafico_radar: Imagem do gráfico radar
            diagnostico: Diagnóstico estratégico
            
        Returns:
            DOCX em bytes
        """
        doc = Document()
        
        # ===== ESTILOS =====
        styles = doc.styles
        
        # ===== CAPA =====
        titulo = doc.add_heading('VinciPitch.AI', 0)
        titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()
        
        subtitulo = doc.add_heading('Relatório de Análise', level=1)
        subtitulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()
        
        empresa_titulo = doc.add_heading(empresa.get('nome', 'N/A'), level=1)
        empresa_titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        info = doc.add_paragraph()
        info.alignment = WD_ALIGN_PARAGRAPH.CENTER
        info.add_run(f"Setor: {empresa.get('setor', 'N/A').title()}\n")
        info.add_run(f"Nota Final: {analise.get('nota_final', 0):.2f}/4.00\n").bold = True
        info.add_run(f"\nGerado em: {datetime.now().strftime('%d/%m/%Y %H:%M')}")
        
        doc.add_page_break()
        
        # ===== SUMÁRIO EXECUTIVO =====
        doc.add_heading('Sumário Executivo', level=1)
        resumo = analise.get('resumo_executivo', 'N/A')
        doc.add_paragraph(resumo)
        
        # ===== CLASSIFICAÇÕES =====
        doc.add_heading('Classificações', level=1)
        
        table = doc.add_table(rows=4, cols=2)
        table.style = 'Table Grid'
        
        classificacoes = [
            ('Potencial', analise.get('classificacao_potencial', 'N/A')),
            ('Risco', analise.get('classificacao_risco', 'N/A')),
            ('Recomendação', analise.get('recomendacao_investimento', 'N/A'))
        ]
        
        # Header
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Critério'
        hdr_cells[1].text = 'Classificação'
        
        for i, (criterio, valor) in enumerate(classificacoes, 1):
            row = table.rows[i].cells
            row[0].text = criterio
            row[1].text = str(valor).title() if valor else 'N/A'
        
        doc.add_paragraph()
        
        # ===== GRÁFICO =====
        if grafico_radar:
            doc.add_heading('Análise por Critérios', level=1)
            
            # Salva temporariamente
            img_stream = io.BytesIO(grafico_radar)
            doc.add_picture(img_stream, width=Inches(5))
            last_paragraph = doc.paragraphs[-1]
            last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_page_break()
        
        # ===== NOTAS DETALHADAS =====
        doc.add_heading('Notas por Critério', level=1)
        
        table = doc.add_table(rows=13, cols=3)
        table.style = 'Table Grid'
        
        # Header
        hdr = table.rows[0].cells
        hdr[0].text = 'Critério'
        hdr[1].text = 'Nota'
        hdr[2].text = 'Justificativa'
        
        for i, (criterio_key, criterio_label) in enumerate(CRITERIOS_LABELS.items(), 1):
            row = table.rows[i].cells
            row[0].text = criterio_label
            row[1].text = f"{analise.get(f'nota_{criterio_key}', 0):.2f}"
            just_key = f"justificativa_{criterio_key.split('_')[0]}"
            row[2].text = str(analise.get(just_key, 'N/A'))[:300]
        
        doc.add_page_break()
        
        # ===== DIAGNÓSTICO =====
        if diagnostico:
            doc.add_heading('Diagnóstico Estratégico', level=1)
            
            doc.add_heading('Pontos Fortes', level=2)
            for ponto in diagnostico.get('pontos_fortes', [])[:5]:
                doc.add_paragraph(ponto, style='List Bullet')
            
            doc.add_heading('Pontos de Atenção', level=2)
            for ponto in diagnostico.get('pontos_fracos', [])[:5]:
                doc.add_paragraph(ponto, style='List Bullet')
            
            doc.add_heading('Oportunidades', level=2)
            for ponto in diagnostico.get('oportunidades', [])[:5]:
                doc.add_paragraph(ponto, style='List Bullet')
            
            doc.add_heading('Ameaças', level=2)
            for ponto in diagnostico.get('ameacas', [])[:5]:
                doc.add_paragraph(ponto, style='List Bullet')
            
            doc.add_heading('Recomendações', level=2)
            for rec in diagnostico.get('recomendacoes', [])[:5]:
                doc.add_paragraph(rec, style='List Bullet')
        
        # Salva em bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        return buffer.read()
    
    # =========================================================================
    # EXPORTAR RANKING COMPLETO
    # =========================================================================
    
    async def exportar_ranking_pdf(
        self,
        ranking: List[Dict[str, Any]],
        estatisticas: Dict[str, Any],
        grafico_barras: Optional[bytes] = None,
        titulo: str = "Ranking VinciPitch.AI"
    ) -> bytes:
        """
        Exporta ranking completo para PDF.
        
        Args:
            ranking: Lista ordenada de empresas com notas
            estatisticas: Estatísticas gerais
            grafico_barras: Gráfico de barras do ranking
            titulo: Título do relatório
            
        Returns:
            PDF em bytes
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=1.5*cm, leftMargin=1.5*cm)
        
        story = []
        
        # Título
        story.append(Paragraph(titulo, self.styles['TituloPrincipal']))
        story.append(Paragraph(f"Gerado em: {datetime.now().strftime('%d/%m/%Y %H:%M')}", self.styles['Normal']))
        story.append(Spacer(1, 1*cm))
        
        # Estatísticas
        story.append(Paragraph("Resumo Estatístico", self.styles['Subtitulo']))
        
        stats_data = [
            ['Métrica', 'Valor'],
            ['Total de Empresas', str(estatisticas.get('total', 0))],
            ['Média Geral', f"{estatisticas.get('media', 0):.2f}"],
            ['Mediana', f"{estatisticas.get('mediana', 0):.2f}"],
            ['Desvio Padrão', f"{estatisticas.get('desvio_padrao', 0):.2f}"]
        ]
        
        table = Table(stats_data, colWidths=[6*cm, 4*cm])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4A148C')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E2E8F0'))
        ]))
        story.append(table)
        story.append(Spacer(1, 1*cm))
        
        # Gráfico
        if grafico_barras:
            story.append(Paragraph("Visualização do Ranking", self.styles['Subtitulo']))
            img = Image(io.BytesIO(grafico_barras), width=16*cm, height=10*cm)
            story.append(img)
            story.append(PageBreak())
        
        # Tabela de ranking
        story.append(Paragraph("Ranking Detalhado", self.styles['Subtitulo']))
        
        ranking_data = [['#', 'Empresa', 'Setor', 'Nota', 'Potencial']]
        for item in ranking[:50]:  # Top 50
            ranking_data.append([
                str(item.get('posicao', '-')),
                item.get('nome', 'N/A')[:25],
                item.get('setor', 'N/A')[:10],
                f"{item.get('nota_final', 0):.2f}",
                item.get('classificacao_potencial', 'N/A')[:10]
            ])
        
        table = Table(ranking_data, colWidths=[1*cm, 5*cm, 3*cm, 2*cm, 3*cm])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4A148C')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E2E8F0')),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8FAFC')])
        ]))
        story.append(table)
        
        doc.build(story)
        buffer.seek(0)
        
        return buffer.read()

    # ═══════════════════════════════════════════════════════════════════════════════
    # EXPORTAÇÃO DE COMPARAÇÃO
    # ═══════════════════════════════════════════════════════════════════════════════

    async def exportar_comparacao_pdf(
        self,
        empresa_a: Dict,
        empresa_b: Dict,
        analise_a: Dict,
        analise_b: Dict,
        grafico_comparativo: bytes = None
    ) -> bytes:
        """Exporta comparação entre duas empresas em PDF."""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm)
        story = []
        
        # Título
        story.append(Paragraph("Comparativo de Startups", self.styles['TituloCapa']))
        story.append(Spacer(1, 0.5*cm))
        story.append(Paragraph(
            f"{empresa_a.get('nome', 'Empresa A')} vs {empresa_b.get('nome', 'Empresa B')}",
            self.styles['SubtituloCapa']
        ))
        story.append(Spacer(1, 0.3*cm))
        story.append(Paragraph(
            f"Gerado em {datetime.now().strftime('%d/%m/%Y às %H:%M')}",
            self.styles['TextoPequeno']
        ))
        story.append(self._criar_linha_decorativa())
        story.append(Spacer(1, 1*cm))
        
        # Resultado geral
        nota_a = analise_a.get('nota_final', 0) or 0
        nota_b = analise_b.get('nota_final', 0) or 0
        
        vencedor = "Empate Técnico"
        if nota_a > nota_b:
            vencedor = empresa_a.get('nome', 'Empresa A')
        elif nota_b > nota_a:
            vencedor = empresa_b.get('nome', 'Empresa B')
        
        story.append(Paragraph("Resultado", self.styles['TituloSecao']))
        
        resultado_data = [
            ['', empresa_a.get('nome', 'A')[:20], empresa_b.get('nome', 'B')[:20]],
            ['Nota Final', f"{nota_a:.2f}", f"{nota_b:.2f}"],
            ['Percentual', f"{nota_a * 25:.0f}%", f"{nota_b * 25:.0f}%"],
            ['Vencedor', vencedor if nota_a >= nota_b else '', vencedor if nota_b > nota_a else '']
        ]
        
        table = Table(resultado_data, colWidths=[4*cm, 5*cm, 5*cm])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), CORES['roxo']),
            ('TEXTCOLOR', (0, 0), (-1, 0), CORES['branco']),
            ('BACKGROUND', (0, 1), (0, -1), CORES['fundo_claro']),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTNAME', (0, 1), (0, -1), 'Helvetica-Bold'),
            ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
            ('GRID', (0, 0), (-1, -1), 0.5, CORES['borda']),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [CORES['branco'], CORES['fundo_claro']])
        ]))
        table.hAlign = 'CENTER'
        story.append(table)
        story.append(Spacer(1, 1*cm))
        
        # Gráfico comparativo
        if grafico_comparativo:
            story.append(Paragraph("Comparativo por Critério", self.styles['TituloSecao']))
            img = Image(io.BytesIO(grafico_comparativo), width=15*cm, height=10*cm)
            img.hAlign = 'CENTER'
            story.append(img)
            story.append(Spacer(1, 1*cm))
        
        # Tabela detalhada por critério
        story.append(Paragraph("Detalhamento por Critério", self.styles['TituloSecao']))
        
        criterios_data = [['Critério', empresa_a.get('nome', 'A')[:15], empresa_b.get('nome', 'B')[:15], 'Diferença', 'Vencedor']]
        
        for key, label in CRITERIOS_LABELS.items():
            nota_crit_a = analise_a.get(f'nota_{key}', 0) or 0
            nota_crit_b = analise_b.get(f'nota_{key}', 0) or 0
            diferenca = nota_crit_a - nota_crit_b
            
            venc = '-'
            if nota_crit_a > nota_crit_b:
                venc = 'A'
            elif nota_crit_b > nota_crit_a:
                venc = 'B'
            
            criterios_data.append([
                label[:20],
                f"{nota_crit_a:.1f}",
                f"{nota_crit_b:.1f}",
                f"{diferenca:+.1f}",
                venc
            ])
        
        table = Table(criterios_data, colWidths=[5*cm, 2.5*cm, 2.5*cm, 2*cm, 2*cm])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), CORES['roxo']),
            ('TEXTCOLOR', (0, 0), (-1, 0), CORES['branco']),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
            ('GRID', (0, 0), (-1, -1), 0.5, CORES['borda']),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [CORES['branco'], CORES['fundo_claro']])
        ]))
        table.hAlign = 'CENTER'
        story.append(table)
        
        # Rodapé
        story.append(Spacer(1, 2*cm))
        story.append(self._criar_linha_decorativa(largura=10*cm))
        story.append(Paragraph(
            "Relatório gerado por VinciPitch.AI",
            ParagraphStyle('Rodape', alignment=TA_CENTER, textColor=CORES['cinza_claro'], fontSize=9)
        ))
        
        doc.build(story)
        buffer.seek(0)
        return buffer.read()

    async def exportar_comparacao_docx(
        self,
        empresa_a: Dict,
        empresa_b: Dict,
        analise_a: Dict,
        analise_b: Dict,
        grafico_comparativo: bytes = None
    ) -> bytes:
        """Exporta comparação entre duas empresas em Word (DOCX)."""
        doc = Document()
        
        # Título
        titulo = doc.add_heading('Comparativo de Startups', level=0)
        titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        subtitulo = doc.add_paragraph(f"{empresa_a.get('nome', 'Empresa A')} vs {empresa_b.get('nome', 'Empresa B')}")
        subtitulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        data_p = doc.add_paragraph(f"Gerado em {datetime.now().strftime('%d/%m/%Y às %H:%M')}")
        data_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()
        
        # Resultado
        nota_a = analise_a.get('nota_final', 0) or 0
        nota_b = analise_b.get('nota_final', 0) or 0
        
        vencedor = "Empate Técnico"
        if nota_a > nota_b:
            vencedor = empresa_a.get('nome', 'Empresa A')
        elif nota_b > nota_a:
            vencedor = empresa_b.get('nome', 'Empresa B')
        
        doc.add_heading('Resultado Geral', level=1)
        
        resultado_table = doc.add_table(rows=4, cols=3)
        resultado_table.style = 'Table Grid'
        
        # Cabeçalho
        resultado_table.cell(0, 0).text = ''
        resultado_table.cell(0, 1).text = empresa_a.get('nome', 'Empresa A')[:20]
        resultado_table.cell(0, 2).text = empresa_b.get('nome', 'Empresa B')[:20]
        
        # Dados
        resultado_table.cell(1, 0).text = 'Nota Final'
        resultado_table.cell(1, 1).text = f"{nota_a:.2f}"
        resultado_table.cell(1, 2).text = f"{nota_b:.2f}"
        
        resultado_table.cell(2, 0).text = 'Percentual'
        resultado_table.cell(2, 1).text = f"{nota_a * 25:.0f}%"
        resultado_table.cell(2, 2).text = f"{nota_b * 25:.0f}%"
        
        resultado_table.cell(3, 0).text = 'Vencedor'
        resultado_table.cell(3, 1).text = '✓' if nota_a >= nota_b and nota_a != nota_b else ''
        resultado_table.cell(3, 2).text = '✓' if nota_b > nota_a else ''
        
        doc.add_paragraph()
        
        # Gráfico
        if grafico_comparativo:
            doc.add_heading('Comparativo Visual', level=1)
            img_stream = io.BytesIO(grafico_comparativo)
            doc.add_picture(img_stream, width=Inches(6))
            doc.add_paragraph()
        
        # Tabela detalhada
        doc.add_heading('Detalhamento por Critério', level=1)
        
        criterios_table = doc.add_table(rows=len(CRITERIOS_LABELS) + 1, cols=5)
        criterios_table.style = 'Table Grid'
        
        # Cabeçalho
        criterios_table.cell(0, 0).text = 'Critério'
        criterios_table.cell(0, 1).text = empresa_a.get('nome', 'A')[:15]
        criterios_table.cell(0, 2).text = empresa_b.get('nome', 'B')[:15]
        criterios_table.cell(0, 3).text = 'Diferença'
        criterios_table.cell(0, 4).text = 'Vencedor'
        
        # Dados
        for i, (key, label) in enumerate(CRITERIOS_LABELS.items(), 1):
            nota_crit_a = analise_a.get(f'nota_{key}', 0) or 0
            nota_crit_b = analise_b.get(f'nota_{key}', 0) or 0
            diferenca = nota_crit_a - nota_crit_b
            
            venc = '-'
            if nota_crit_a > nota_crit_b:
                venc = empresa_a.get('nome', 'A')[:10]
            elif nota_crit_b > nota_crit_a:
                venc = empresa_b.get('nome', 'B')[:10]
            
            criterios_table.cell(i, 0).text = label
            criterios_table.cell(i, 1).text = f"{nota_crit_a:.1f}"
            criterios_table.cell(i, 2).text = f"{nota_crit_b:.1f}"
            criterios_table.cell(i, 3).text = f"{diferenca:+.1f}"
            criterios_table.cell(i, 4).text = venc
        
        doc.add_paragraph()
        
        # Rodapé
        rodape = doc.add_paragraph('Relatório gerado por VinciPitch.AI')
        rodape.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Salva
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.read()